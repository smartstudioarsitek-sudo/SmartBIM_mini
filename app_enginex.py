import streamlit as st
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import json
from PIL import Image
import PyPDF2
import io
import docx
import zipfile
from pptx import Presentation
import re

# --- KONFIGURASI HALAMAN ---
st.set_page_config(page_title="ENGINEX Ultimate", page_icon="🏗️", layout="wide")

# --- CSS BIAR TAMPILAN GAGAH ---
st.markdown("""
<style>
    .main-header {font-size: 30px; font-weight: bold; color: #1E3A8A; margin-bottom: 10px;}
    [data-testid="stSidebar"] {background-color: #f8f9fa;}
    .stChatInput textarea {font-size: 16px !important;}
    
    /* Efek Avatar */
    .stChatMessage .avatar {background-color: #1E3A8A; color: white;}
    
    /* Tombol Download Custom */
    .stDownloadButton button {
        width: 100%;
        border-radius: 8px;
        font-weight: bold;
    }
    
    /* Highlight untuk Mode Auto-Pilot */
    .auto-pilot-msg {
        background-color: #e0f7fa;
        border-left: 5px solid #00acc1;
        padding: 10px;
        margin-bottom: 10px;
        border-radius: 5px;
        color: #006064;
        font-weight: bold;
    }
    
    /* Highlight Grafik */
    .plot-container {
        border: 1px solid #ddd;
        border-radius: 10px;
        padding: 10px;
        margin-top: 10px;
        background-color: white;
    }
</style>
""", unsafe_allow_html=True)

# --- INIT SESSION STATE ---
if 'processed_files' not in st.session_state:
    st.session_state.processed_files = set()
if 'current_expert_active' not in st.session_state:
    st.session_state.current_expert_active = "👑 The GEMS Grandmaster"

# ==========================================
# 0. FUNGSI BANTUAN EXPORT & PLOTTING
# ==========================================

def create_docx_from_text(text_content):
    """Mengubah teks chat menjadi file Word (.docx)"""
    try:
        doc = docx.Document()
        doc.add_heading('Laporan Output ENGINEX', 0)
        
        lines = text_content.split('\n')
        for line in lines:
            clean_line = line.strip()
            if clean_line.startswith('## '):
                doc.add_heading(clean_line.replace('## ', ''), level=2)
            elif clean_line.startswith('### '):
                doc.add_heading(clean_line.replace('### ', ''), level=3)
            elif clean_line.startswith('- ') or clean_line.startswith('* '):
                try:
                    doc.add_paragraph(clean_line, style='List Bullet')
                except:
                    doc.add_paragraph(clean_line)
            elif clean_line:
                doc.add_paragraph(clean_line)
                
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio
    except Exception as e:
        return None

def extract_table_to_excel(text_content):
    """Mendeteksi tabel Markdown dalam chat dan mengubahnya ke Excel (.xlsx)"""
    try:
        lines = text_content.split('\n')
        table_data = []
        
        for line in lines:
            stripped = line.strip()
            if "|" in stripped:
                if set(stripped.replace('|', '').replace('-', '').replace(' ', '')) == set():
                    continue
                row_cells = [c.strip() for c in stripped.split('|')]
                if stripped.startswith('|'): row_cells = row_cells[1:]
                if stripped.endswith('|'): row_cells = row_cells[:-1]
                if row_cells:
                    table_data.append(row_cells)
        
        if len(table_data) < 2: return None
            
        headers = table_data[0]
        data_rows = table_data[1:]
        df = pd.DataFrame(data_rows, columns=headers)
        
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False, sheet_name='Data_ENGINEX')
            worksheet = writer.sheets['Data_ENGINEX']
            for i, col in enumerate(df.columns):
                worksheet.set_column(i, i, 20)
        output.seek(0)
        return output
    except Exception as e:
        return None

def execute_generated_code(code_str):
    """
    [ENGINEERING PLOTTER]
    Mengeksekusi string kode Python yang dihasilkan AI untuk membuat grafik.
    """
    try:
        # Create a dictionary for local variables
        local_vars = {
            "pd": pd,
            "np": np,
            "plt": plt,
            "st": st
        }
        
        # Eksekusi kode dalam lingkungan aman
        exec(code_str, {}, local_vars)
        return True
    except Exception as e:
        st.error(f"⚠️ Gagal Render Grafik: {e}")
        return False

# ==========================================
# 1. SETUP API KEY & MODEL (SIDEBAR)
# ==========================================
with st.sidebar:
    st.title("🏗️ ENGINEX ULTIMATE")
    st.caption("v10.0 | Agentic + Plotting Engine")
    
    api_key_input = st.text_input("🔑 API Key:", type="password")
    if api_key_input:
        raw_key = api_key_input
        st.caption("ℹ️ Key Manual Digunakan")
    else:
        raw_key = st.secrets.get("GOOGLE_API_KEY")
    
    if not raw_key:
        st.warning("⚠️ Masukkan API Key Google AI Studio.")
        st.stop()
        
    clean_api_key = raw_key.strip()

try:
    genai.configure(api_key=clean_api_key, transport="rest")
except Exception as e:
    st.error(f"Config Error: {e}")

@st.cache_resource
def get_available_models_from_google(api_key_trigger):
    try:
        model_list = []
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                model_list.append(m.name)
        model_list.sort(key=lambda x: 'pro' not in x) 
        return model_list, None
    except Exception as e:
        return [], str(e)

real_models, error_msg = get_available_models_from_google(clean_api_key)

with st.sidebar:
    if error_msg: st.error(f"❌ Error: {error_msg}"); st.stop()
    if not real_models: st.warning("⚠️ Tidak ada model."); st.stop()

    default_idx = 0
    for i, m in enumerate(real_models):
        if "flash" in m:  
            default_idx = i
            break
            
    selected_model_name = st.selectbox("🧠 Pilih Otak AI:", real_models, index=default_idx)
    
    if "pro" in selected_model_name or "ultra" in selected_model_name:
        st.success(f"⚡ Mode: HIGH REASONING")
    else:
        st.info(f"🚀 Mode: HIGH SPEED")
    st.divider()

# --- KONEKSI DATABASE ---
try:
    from backend_enginex import EnginexBackend
    if 'backend' not in st.session_state:
        st.session_state.backend = EnginexBackend()
    db = st.session_state.backend
except ImportError:
    st.error("⚠️ File 'backend_enginex.py' belum ada!")
    st.stop()

# ==========================================
# 2. SAVE/LOAD & PROYEK
# ==========================================
with st.sidebar:
    with st.expander("💾 Manajemen Data"):
        st.download_button("⬇️ Backup JSON", db.export_data(), "backup.json", mime="application/json")
        uploaded_restore = st.file_uploader("⬆️ Restore", type=["json"])
        if uploaded_restore and st.button("Restore"):
            ok, msg = db.import_data(uploaded_restore)
            if ok: st.success(msg); st.rerun()
            else: st.error(msg)
    
    st.divider()
    existing_projects = db.daftar_proyek()
    mode_proyek = st.radio("Folder Proyek:", ["Proyek Baru", "Buka Lama"], horizontal=True)
    
    if mode_proyek == "Proyek Baru":
        nama_proyek = st.text_input("Nama Proyek:", "DED Irigasi 2026")
    else:
        nama_proyek = st.selectbox("Pilih Proyek:", existing_projects) if existing_projects else "Belum ada"
    st.divider()

# ==========================================
# 3. DEFINISI PERSONA (UPDATED WITH PLOTTING INSTRUCTIONS)
# ==========================================

PLOT_INSTRUCTION = """
[ATURAN PENTING UNTUK VISUALISASI DATA]:
Jika user meminta grafik/diagram/plot:
1. JANGAN HANYA MEMBERIKAN DESKRIPSI.
2. ANDA WAJIB MENULISKAN KODE PYTHON DI DALAM BLOK KODE (```python).
3. Gunakan library `matplotlib.pyplot` (sebagai plt) dan `numpy` (sebagai np).
4. WAJIB: Di akhir kode plotting, gunakan perintah `st.pyplot(plt.gcf())` agar grafik muncul di layar aplikasi Streamlit user.
5. Jangan gunakan `plt.show()`.
6. Berikan judul, label sumbu, dan grid agar grafik terlihat profesional teknik sipil.
"""

gems_persona = {
        "👑 The GEMS Grandmaster": f"""
        ANDA ADALAH "THE GEMS GRANDMASTER" (Omniscient Project Director).
        Anda memiliki 5 "MODUL OTAK":
        1. MODUL DIREKSI & LEGAL (Project Manager)
        2. MODUL HIKMAH & SYARIAH (Ulama Fiqih Bangunan)
        3. MODUL ENGINEERING FISIK (Sipil/SDA/MEP)
        4. MODUL ARSITEKTUR & VISUAL (Konseptor)
        5. MODUL DIGITAL & TOOLS (Coder & Plotter)

        {PLOT_INSTRUCTION}

        INSTRUKSI RESPON:
        1. Analisis Multi-Dimensi (Teknis, Biaya, Hukum, Agama).
        2. Format Profesional (Heading, Bullet points).
        3. Jika perlu hitungan kompleks atau grafik, tuliskan kode Python.
    """,
       "👔 Project Manager (PM)": """
        ANDA ADALAH SENIOR PROJECT DIRECTOR (PMP Certified).
        TUGAS: Keputusan strategis, mitigasi risiko, manajemen stakeholders.
        GAYA: Tegas, Solutif, Strategis.
    """,
    "📝 Drafter Laporan DED": """
        ANDA ADALAH LEAD TECHNICAL WRITER.
        TUGAS: Menyusun Laporan (Pendahuluan, Antara, Akhir) standar PUPR.
    """,
    "⚖️ Ahli Legal & Kontrak": """
        ANDA ADALAH AHLI HUKUM KONSTRUKSI (FIDIC).
        TUGAS: Analisis kontrak, klaim, sengketa, dan regulasi.
    """,
    "🕌 Dewan Syariah": """
        ANDA ADALAH GRAND MUFTI FIQIH BANGUNAN.
        TUGAS: Fatwa arah kiblat, akad syariah, adab membangun.
    """,
    "🌾 Ahli IKSI-PAI": f"""
        ANDA ADALAH PRINCIPAL IRRIGATION ENGINEER (Permen PUPR).
        TUGAS: IKSI, PAI, Audit Irigasi.
        {PLOT_INSTRUCTION}
    """,
    "🌊 Ahli Bangunan Air": f"""
        ANDA ADALAH SENIOR HYDRAULIC ENGINEER.
        TUGAS: Bendung, Bendungan, Pintu Air.
        {PLOT_INSTRUCTION}
    """,
    "🌧️ Ahli Hidrologi": f"""
        ANDA ADALAH SENIOR HYDROLOGIST.
        TUGAS: Analisis Curah Hujan, Banjir Rencana, Debit Andalan.
        {PLOT_INSTRUCTION}
    """,
    "🏖️ Ahli Teknik Pantai": f"""
        ANDA ADALAH COASTAL ENGINEER.
        TUGAS: Breakwater, Pasang Surut, Reklamasi.
        {PLOT_INSTRUCTION}
    """,
    "🏗️ Ahli Struktur (Gedung)": f"""
        ANDA ADALAH PRINCIPAL STRUCTURAL ENGINEER.
        TUGAS: Analisis Gempa, Beton, Baja.
        {PLOT_INSTRUCTION}
    """,
    "🪨 Ahli Geoteknik": f"""
        ANDA ADALAH SENIOR GEOTECHNICAL ENGINEER.
        TUGAS: Pondasi, Stabilitas Lereng, Sondir.
        {PLOT_INSTRUCTION}
    """,
    "🛣️ Ahli Jalan & Jembatan": f"""
        ANDA ADALAH HIGHWAY ENGINEER.
        TUGAS: Geometrik Jalan, Perkerasan, Jembatan.
        {PLOT_INSTRUCTION}
    """,
    "🌍 Ahli Geodesi & GIS": """
        ANDA ADALAH GEOMATICS ENGINEER.
        TUGAS: Survey, Peta Kontur, Cut & Fill.
    """,
    "🏛️ Senior Architect": """
        ANDA ADALAH PRINCIPAL ARCHITECT.
        TUGAS: Desain, Estetika, Fungsi Ruang.
    """,
    "🌳 Landscape Architect": """
        ANDA ADALAH LANDSCAPE ARCHITECT.
        TUGAS: RTH, Taman, Drainase Lingkungan.
    """,
    "🎨 The Visionary Architect": """
        ANDA ADALAH PROMPT ENGINEER & VISUALIZER.
        TUGAS: Membuat "Master Prompt" untuk Image Generator berdasarkan ide user.
    """,
    "🌍 Ahli Planologi": """
        ANDA ADALAH URBAN PLANNER.
        TUGAS: Tata Ruang, Zonasi, Masterplan.
    """,
    "🏭 Ahli Proses Industri": """
        ANDA ADALAH PROCESS ENGINEER.
        TUGAS: PFD, P&ID, Pabrik Kimia.
    """,
    "📜 Ahli AMDAL": """
        ANDA ADALAH KETUA TIM AMDAL.
        TUGAS: Dokumen Lingkungan, Mitigasi Dampak.
    """,
    "♻️ Ahli Teknik Lingkungan": """
        ANDA ADALAH SANITARY ENGINEER.
        TUGAS: IPAL, WTP, Persampahan.
    """,
    "⛑️ Ahli K3 Konstruksi": """
        ANDA ADALAH SAFETY MANAGER.
        TUGAS: CSMS, IBPRP, SMKK.
    """,
    "💻 Lead Engineering Developer": f"""
        ANDA ADALAH LEAD FULL-STACK ENGINEER.
        KEAHLIAN: Python, Streamlit, Plotting Data.
        {PLOT_INSTRUCTION}
    """,
    "📐 CAD & BIM Automator": """
        ANDA ADALAH BIM MANAGER.
        TUGAS: Revit API, Dynamo, Scripting.
    """,
    "🖥️ Instruktur Software": """
        ANDA ADALAH TRAINER SOFTWARE.
        TUGAS: Tutorial Civil 3D, HEC-RAS, SAP2000.
    """,
    "💰 Ahli Estimator (RAB)": """
        ANDA ADALAH QUANTITY SURVEYOR.
        TUGAS: RAB, AHSP, BoQ.
    """,
    "💵 Ahli Keuangan Proyek": f"""
        ANDA ADALAH FINANCE MANAGER.
        TUGAS: NPV, IRR, Cashflow.
        {PLOT_INSTRUCTION}
    """,
    "📜 Ahli Perizinan": """
        ANDA ADALAH KONSULTAN PERIZINAN.
        TUGAS: PBG, SLF, KRK.
    """,
    "🤖 The Enginex Architect": """
        ANDA ADALAH SYSTEM ADMINISTRATOR APLIKASI INI.
    """,
}

# ==========================================
# 4. FUNGSI AUTO-ROUTER
# ==========================================
def get_auto_pilot_decision(user_query, model_api_key):
    try:
        router_model = genai.GenerativeModel("gemini-1.5-flash")
        list_ahli = list(gems_persona.keys())
        router_prompt = f"""
        Pilih SATU ahli dari daftar berikut untuk menjawab pertanyaan: "{user_query}"
        Daftar: {list_ahli}
        Output: HANYA nama ahli persis. Jika ragu, pilih '👑 The GEMS Grandmaster'.
        """
        response = router_model.generate_content(router_prompt)
        suggested = response.text.strip()
        if suggested in list_ahli: return suggested
        return "👑 The GEMS Grandmaster"
    except:
        return "👑 The GEMS Grandmaster"

# ==========================================
# 5. SIDEBAR BAWAH & FILE UPLOAD
# ==========================================
with st.sidebar:
    st.markdown("### 👷 Tim Ahli")
    use_auto_pilot = st.checkbox("🤖 Auto-Pilot", value=True)
    manual_selection = st.selectbox("Pilih Manual:", list(gems_persona.keys()), disabled=use_auto_pilot)
    
    if not use_auto_pilot:
        st.session_state.current_expert_active = manual_selection
    
    st.markdown("---")
    st.markdown("### 📂 Upload Data")
    uploaded_files = st.file_uploader(
        "File:", 
        type=["png", "jpg", "jpeg", "pdf", "docx", "doc", "xlsx", "xls", "pptx", "zip", "dwg", "kml", "kmz", "geojson", "gpx", "py"], 
        accept_multiple_files=True
    )
    
    if uploaded_files: st.info(f"📎 {len(uploaded_files)} File")
    
    st.divider()
    if st.button("🧹 Reset Chat"):
        db.clear_chat(nama_proyek, st.session_state.current_expert_active)
        st.session_state.processed_files.clear()
        st.rerun()

# ==========================================
# 6. FUNGSI BACA FILE
# ==========================================
def process_uploaded_file(uploaded_file):
    if uploaded_file is None: return None, None
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    try:
        if file_type in ['png', 'jpg', 'jpeg']:
            return "image", Image.open(uploaded_file)
        elif file_type == 'pdf':
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages: 
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
            return "text", text
        elif file_type == 'docx':
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
            return "text", text
        elif file_type == 'doc':
            try:
                raw_data = uploaded_file.getvalue()
                text = "".join([chr(b) for b in raw_data if 32 <= b <= 126 or b in [10, 13]])
                return "text", f"[RAW READ .DOC]\n{text}"
            except Exception as e:
                return "error", str(e)
        elif file_type in ['xlsx', 'xls']:
            try:
                df = pd.read_excel(uploaded_file)
                return "text", f"[PREVIEW EXCEL]\n{df.head(50).to_csv(index=False)}"
            except Exception as e:
                return "error", str(e)
        elif file_type == 'pptx':
            prs = Presentation(uploaded_file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text.append(shape.text)
            return "text", "\n".join(text)
        elif file_type == 'py':
            return "text", uploaded_file.getvalue().decode("utf-8")
        elif file_type in ['kml', 'geojson', 'gpx']:
            return "text", uploaded_file.getvalue().decode("utf-8")
        elif file_type == 'kmz':
            with zipfile.ZipFile(uploaded_file, "r") as z:
                kml = [n for n in z.namelist() if n.endswith(".kml")][0]
                with z.open(kml) as f: return "text", f.read().decode("utf-8")
        elif file_type == 'zip':
            with zipfile.ZipFile(uploaded_file, "r") as z:
                return "text", f"ZIP Content:\n{', '.join(z.namelist())}"
    except Exception as e: 
        return "error", str(e)
    return "error", "Format tidak didukung"

# ==========================================
# 7. MAIN CHAT AREA
# ==========================================
st.markdown(f'<div class="main-header">{nama_proyek}</div>', unsafe_allow_html=True)

current_expert = st.session_state.current_expert_active
st.caption(f"Status: **Connected** | Expert: **{current_expert}**")

# Display History
history = db.get_chat_history(nama_proyek, current_expert)
for chat in history:
    with st.chat_message(chat['role']):
        st.markdown(chat['content'])
        # [NEW] Check for Plot Code in History to Re-render? 
        # (Optional: Usually history is text only, re-rendering might be heavy. Skipped for performance)

prompt = st.chat_input(f"Tanya sesuatu ke {current_expert}...")

if prompt:
    # --- AUTO PILOT ---
    detected_expert = current_expert
    if use_auto_pilot:
        with st.status("🧠 Menganalisis konteks...", expanded=True) as status:
            detected_expert = get_auto_pilot_decision(prompt, clean_api_key)
            status.write(f"Ahli yang relevan: **{detected_expert}**")
            st.session_state.current_expert_active = detected_expert
            st.markdown(f'<div class="auto-pilot-msg">🤖 Auto-Pilot: Mengalihkan ke {detected_expert}</div>', unsafe_allow_html=True)
    
    final_expert_name = detected_expert

    # --- SAVE USER CHAT ---
    db.simpan_chat(nama_proyek, final_expert_name, "user", prompt)
    with st.chat_message("user"):
        st.markdown(prompt)
    
    # --- PREPARE CONTEXT ---
    content_to_send = [prompt]
    if uploaded_files:
        for upl_file in uploaded_files:
            if upl_file.name not in st.session_state.processed_files:
                ftype, fcontent = process_uploaded_file(upl_file)
                if ftype == "image":
                    with st.chat_message("user"): st.image(upl_file, width=200)
                    content_to_send.append(fcontent)
                elif ftype == "text":
                    with st.chat_message("user"): st.caption(f"📄 Data: {upl_file.name}")
                    content_to_send[0] += f"\n\n--- FILE: {upl_file.name} ---\n{fcontent}\n------\n"
                st.session_state.processed_files.add(upl_file.name)

    # --- GENERATE AI RESPONSE ---
    with st.chat_message("assistant"):
        with st.spinner(f"{final_expert_name.split(' ')[1]} sedang berpikir & plot grafik..."):
            try:
                safety = {
                    HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
                }

                model = genai.GenerativeModel(
                    model_name=selected_model_name,
                    system_instruction=gems_persona[final_expert_name], 
                    safety_settings=safety
                )
                
                # Context History
                current_history = db.get_chat_history(nama_proyek, final_expert_name)
                hist_formatted = []
                for h in current_history:
                    if h['content'] != prompt:
                        role_api = "user" if h['role']=="user" else "model"
                        hist_formatted.append({"role": role_api, "parts": [h['content']]})
                
                chat_session = model.start_chat(history=hist_formatted)
                response_stream = chat_session.send_message(content_to_send, stream=True)
                
                full_response_text = ""
                placeholder = st.empty()
                
                for chunk in response_stream:
                    if chunk.text:
                        full_response_text += chunk.text
                        placeholder.markdown(full_response_text + "▌")
                
                placeholder.markdown(full_response_text)
                db.simpan_chat(nama_proyek, final_expert_name, "assistant", full_response_text)
                
                # ==================================================
                # [NEW FEATURE] ENGINEERING PLOTTER EXECUTION
                # ==================================================
                # Regex untuk mencari blok kode Python
                code_blocks = re.findall(r"```python(.*?)```", full_response_text, re.DOTALL)
                
                for code in code_blocks:
                    if "plt." in code or "matplotlib" in code:
                        st.markdown("### 📊 Engineering Plotter Output:")
                        with st.container():
                            # Jalankan kode plotting
                            success = execute_generated_code(code)
                            if success:
                                st.caption("✅ Grafik berhasil di-render dari kode Python.")
                            # Bersihkan plot agar tidak tumpah ke chat berikutnya
                            plt.clf()

                # ==================================================
                # DOWNLOAD BUTTONS
                # ==================================================
                st.markdown("---")
                col1, col2 = st.columns(2)
                
                docx_file = create_docx_from_text(full_response_text)
                if docx_file:
                    col1.download_button("📄 Download Laporan (.docx)", docx_file, f"Laporan_{final_expert_name[:5]}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                
                xlsx_file = extract_table_to_excel(full_response_text)
                if xlsx_file:
                    col2.download_button("📊 Download Tabel (.xlsx)", xlsx_file, f"Data_{final_expert_name[:5]}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                
            except Exception as e:
                st.error(f"⚠️ Error: {e}")
